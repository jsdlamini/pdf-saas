#!/usr/bin/env python3
"""PDF → PPTX using typography inference.

Each PDF page becomes a slide. The largest-font text block is treated as the
slide title; the remaining blocks become bulleted body text. Embedded images
are extracted and placed on the right of the slide.

Usage: python3 pdf2powerpoint.py <input.pdf> <output.pptx>
Prints "SLIDES <n> IMAGES <m>" on success for the route's engine stats.
"""
import os
import sys

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt


def main() -> None:
    if len(sys.argv) < 3:
        print("usage: pdf2powerpoint.py <input.pdf> <output.pptx>", file=sys.stderr)
        sys.exit(2)

    input_path = sys.argv[1]
    output_path = sys.argv[2]
    scratch_dir = os.path.dirname(output_path) or "."

    doc = fitz.open(input_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    image_count = 0

    for page_index in range(len(doc)):
        page = doc[page_index]
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        # Collect text lines with their largest span size + bold flag.
        lines = []
        text_dict = page.get_text("dict")
        for block in text_dict.get("blocks", []):
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                spans = line.get("spans", [])
                if not spans:
                    continue
                text = "".join(span["text"] for span in spans).strip()
                if not text:
                    continue
                size = max((span.get("size") or 12) for span in spans)
                bold = any("Bold" in (span.get("font") or "") for span in spans)
                lines.append((size, bold, text))

        # Title = largest text; body = the rest, deduplicated.
        lines.sort(key=lambda item: (-item[0], -int(item[1])))
        if lines:
            size, _bold, title_text = lines[0]
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12.1), Inches(1.1))
            title_frame = title_box.text_frame
            title_frame.text = title_text
            for paragraph in title_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(min(max(size, 24), 40))
                    run.font.bold = True

            body = lines[1:14]
            if body:
                body_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5.2))
                body_frame = body_box.text_frame
                seen = set()
                first = True
                for size_b, _bold_b, text_b in body:
                    if text_b in seen:
                        continue
                    seen.add(text_b)
                    paragraph = body_frame.paragraphs[0] if first else body_frame.add_paragraph()
                    first = False
                    paragraph.text = text_b if text_b.startswith("•") else f"• {text_b}"
                    for run in paragraph.runs:
                        run.font.size = Pt(min(max(size_b, 12), 20))

        # Embedded images placed on the right.
        try:
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                try:
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n - pix.alpha >= 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    img_path = os.path.join(scratch_dir, f"pg{page_index}_img{xref}.png")
                    pix.save(img_path)
                    slide.shapes.add_picture(img_path, Inches(9.2), Inches(1.0), width=Inches(3.6))
                    image_count += 1
                except Exception:
                    continue
        except Exception:
            pass

    prs.save(output_path)
    print(f"SLIDES {len(doc)} IMAGES {image_count}")


if __name__ == "__main__":
    main()
