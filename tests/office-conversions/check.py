#!/usr/bin/env python3
"""Structure check for the PDF -> Excel and PDF -> PowerPoint converters.

Generates a small fixture PDF containing a real table and an image, runs
scripts/pdf2excel.py and scripts/pdf2powerpoint.py, then asserts the output
carries real structure:

  * XLSX has a worksheet whose cells span multiple columns (a real table),
  * PPTX has at least one slide with text (title/bullets) and an image.

Exits non-zero if any assertion fails, so it can gate CI.

Usage (from repo root):
    python3 tests/office-conversions/check.py
"""

import subprocess
import sys
import tempfile
from pathlib import Path

import fitz  # PyMuPDF
from openpyxl import load_workbook
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent.parent
EXCEL = ROOT / "scripts" / "pdf2excel.py"
POWERPOINT = ROOT / "scripts" / "pdf2powerpoint.py"


def make_table_pdf(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page(width=400, height=320)
    margin = 40
    cols = 3
    rows = 3
    cell_w = (400 - 2 * margin) / cols
    cell_h = (320 - 2 * margin) / rows
    for i in range(cols + 1):
        x = margin + i * cell_w
        page.draw_line(fitz.Point(x, margin), fitz.Point(x, 320 - margin))
    for j in range(rows + 1):
        y = margin + j * cell_h
        page.draw_line(fitz.Point(margin, y), fitz.Point(400 - margin, y))

    headers = ["Name", "Qty", "Price"]
    data = [["Apples", "10", "1.50"], ["Pears", "5", "2.00"]]
    for j in range(rows):
        for i in range(cols):
            text = headers[i] if j == 0 else data[j - 1][i]
            page.insert_text(
                fitz.Point(margin + i * cell_w + 8, margin + j * cell_h + cell_h / 2 + 4),
                text,
                fontsize=11,
            )
    doc.save(str(path))
    doc.close()


def make_slide_pdf(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page(width=960, height=540)
    page.insert_text(fitz.Point(60, 90), "Quarterly Report", fontsize=34)
    page.insert_text(fitz.Point(60, 160), "Revenue grew 12%", fontsize=16)
    page.insert_text(fitz.Point(60, 190), "Costs fell 4%", fontsize=16)
    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 120, 80))
    pix.clear_with(200)
    page.insert_image(fitz.Rect(700, 120, 900, 240), stream=pix.tobytes("png"))
    doc.save(str(path))
    doc.close()


def run(script: Path, input_path: Path, output_path: Path):
    proc = subprocess.run(
        [sys.executable, str(script), str(input_path), str(output_path)],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    return proc.returncode, proc.stderr


def main() -> int:
    failures = []

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        table_pdf = tmp_dir / "table.pdf"
        slide_pdf = tmp_dir / "slide.pdf"
        make_table_pdf(table_pdf)
        make_slide_pdf(slide_pdf)

        # PDF -> Excel
        xlsx = tmp_dir / "out.xlsx"
        code, err = run(EXCEL, table_pdf, xlsx)
        if code != 0:
            failures.append(f"pdf2excel exited {code}: {err.strip()}")
        else:
            wb = load_workbook(xlsx)
            sheets = wb.sheetnames
            max_cols = 0
            for sheet_name in sheets:
                ws = wb[sheet_name]
                max_cols = max(max_cols, ws.max_column or 0)
            print(f"pdf2excel: sheets={sheets} max_cols={max_cols}")
            if max_cols < 2:
                failures.append(f"pdf2excel: expected multiple columns, got max_cols={max_cols}")

        # PDF -> PowerPoint
        pptx = tmp_dir / "out.pptx"
        code, err = run(POWERPOINT, slide_pdf, pptx)
        if code != 0:
            failures.append(f"pdf2powerpoint exited {code}: {err.strip()}")
        else:
            prs = Presentation(pptx)
            slides = len(prs.slides)
            text_chars = 0
            pictures = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_chars += len(shape.text_frame.text)
                    if shape.shape_type == 13:  # PICTURE
                        pictures += 1
            print(f"pdf2powerpoint: slides={slides} text_chars={text_chars} pictures={pictures}")
            if slides < 1:
                failures.append("pdf2powerpoint: no slides")
            if text_chars < 10:
                failures.append(f"pdf2powerpoint: expected slide text, got {text_chars} chars")
            if pictures < 1:
                failures.append("pdf2powerpoint: expected an embedded image")

    if failures:
        print("\nFAILURES:")
        for f in failures:
            print(f"  - {f}")
        return 1
    print("\nAll office-conversion checks passed.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
